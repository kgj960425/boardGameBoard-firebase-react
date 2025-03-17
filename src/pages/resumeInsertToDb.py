import os

import psycopg2

from psycopg2 import sql

from docx import Document

from pptx import Presentation



# DB 접속 정보

DB_HOST = 'localhost'

DB_PORT = 5432

DB_NAME = 'postgres'

DB_USER = 'postgres'

DB_PASSWORD = 'zetta@0123'



# 디렉토리 지정

FILE_DIRECTORY = r"C:\Download\temp-file"



# 처리할 항목들: key는 표 내 대상 텍스트, value는 DB 컬럼명

# docx, pptx 파일의 항목명이 달라 추가 테스트시 수정 필요

FIELDS = {

    "성명": "name_desc",

    "소속": "affiliation_desc",

    "직위": "position_desc",

    # "자격증": "certificate_desc",

    "종목 및 등급": "certificate_desc",

    "기술등급": "tech_level_desc",

    "나이": "age_desc",

    "최종학력": "last_education_desc",

    "학위": "degree_desc",

    # "전공": "major_desc",

    "학과(전공)": "major_desc",

    "경력": "work_experience_desc"

}



# DB 커넥션

def get_db_connection():

    try:

        conn = psycopg2.connect(

            host=DB_HOST,

            port=DB_PORT,

            dbname=DB_NAME,

            user=DB_USER,

            password=DB_PASSWORD

        )

        return conn

    except psycopg2.Error as e:

        print("데이터베이스 연결 실패:", e)

        return None



# DB insert 함수

def insert_info(data: dict):

    conn = get_db_connection()

    if conn is None:

        return



    try:

        with conn.cursor() as cur:

            columns = ','.join(data.keys())

            placeholders = ','.join(['%s'] * len(data))

            query = sql.SQL(

                f"""

                INSERT INTO tb_resumes_m ({columns})

                VALUES ({placeholders})

                """

            )

            cur.execute(query, tuple(data.values()))

        conn.commit()

        print("DB 삽입 성공:", data)

    except psycopg2.Error as e:

        print("DB 삽입 실패:", e)

        conn.rollback()

    finally:

        conn.close()



# key 우측 텍스트를 value로

def get_text_right(table, target_text):



    for row in table.rows:

        cells = row.cells

        for idx, cell in enumerate(cells):

            if cell.text.strip() == target_text:

                if idx + 1 < len(cells):

                    return cells[idx + 1].text.strip()

    return ""



# key 하단 텍스트 value로

def get_text_below(table, target_text):



    rows = table.rows

    for row_idx, row in enumerate(rows):

        cells = row.cells

        for col_idx, cell in enumerate(cells):

            if cell.text.strip() == target_text:

                if row_idx + 1 < len(rows):

                    return rows[row_idx + 1].cells[col_idx].text.strip()

    return ""



# below로만 처리하기 위한 함수 (docx 테스트 목적)

def extract_table_data_below(table):

    data = {}

    for key, db_field in FIELDS.items():

        value = get_text_below(table, key)

        if value:

            data[db_field] = value

    return data



# right 먼저 호출 후 없으면 below 호출

def get_value_from_table(table, target_text):

    value = get_text_right(table, target_text)

    if not value:

        value = get_text_below(table, target_text)

    return value



# 표에서 지정한 필드 항목 값 추출 -> 딕셔너리로 반환, 항목 중복시 초기값으로

def extract_table_data(table):



    data = {}

    for key, db_field in FIELDS.items():

        value = get_value_from_table(table, key)

        if value:

            data[db_field] = value

    return data



# word 파일 처리

def process_word_file(file_path):

    try:

        doc = Document(file_path)

    except Exception as e:

        print(f"Word 파일 처리 실패 ({file_path}):", e)

        return



    data = {}

    # 문서 내의 모든 표를 순회하면서 각 표에서 데이터를 추출

    for table in doc.tables:

        table_data = extract_table_data_below(table)

        for key, value in table_data.items():

            if key not in data or not data[key]:

                data[key] = value

        if any(data.get(db_field) for db_field in FIELDS.values()):

            break



    if any(data.get(db_field) for db_field in FIELDS.values()):

        insert_info(data)

    else:

        print(f"파일 내에 필수 항목 내용이 없습니다: {file_path}")



# ppt 파일 처리

def process_ppt_file(file_path):

    try:

        prs = Presentation(file_path)

    except Exception as e:

        print(f"PPT 파일 처리 실패 ({file_path}):", e)

        return



    data = {}

    # 각 슬라이드의 테이블을 순회

    for slide in prs.slides:

        for shape in slide.shapes:

            if not shape.has_table:

                continue

            table = shape.table

            table_data = extract_table_data(table)

            for key, value in table_data.items():

                if key not in data or not data[key]:

                    data[key] = value

            if any(data.get(db_field) for db_field in FIELDS.values()):

                break

        if any(data.get(db_field) for db_field in FIELDS.values()):

            break



    if any(data.get(db_field) for db_field in FIELDS.values()):

        insert_info(data)

    else:

        print(f"파일 내에 필수 항목 내용이 없습니다: {file_path}")




# 특정 디렉토리 내 파일 확장자별 처리

def process_files_in_directory(directory):



    for filename in os.listdir(directory):

        file_path = os.path.join(directory, filename)

        if os.path.isfile(file_path):

            ext = os.path.splitext(filename)[1].lower()

            if ext == '.docx':

                print(f"Word 파일 처리 시작: {file_path}")

                process_word_file(file_path)

            # elif ext == '.pptx':

            #     print(f"PPT 파일 처리 시작: {file_path}")

                process_ppt_file(file_path)

            else:

                print(f"지원되지 않는 확장자: {file_path}")




if __name__ == "__main__":

    process_files_in_directory(FILE_DIRECTORY)